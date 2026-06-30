"""
Extractor de PPTX (Fase 2).

Superficies de ocultamiento propias de las presentaciones:
  - Shapes posicionados fuera del area del slide   -> OFF_PAGE
  - Slides ocultos (atributo show="0")             -> NOTES/structural
  - Notas del orador                               -> NOTES
  - Texto casi blanco                              -> NEAR_WHITE
  - Fuente diminuta                                -> TINY_FONT
  - Texto alternativo de imagenes/shapes           -> ALT_TEXT

Usamos python-pptx. Las posiciones vienen en EMU (English Metric Units):
1 pulgada = 914400 EMU. El slide mide prs.slide_width x prs.slide_height.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.util import Emu

from veilscan.core.models import Finding, HideReason, Severity, Technique, TextSpan
from veilscan.extractors.base import BaseExtractor, ExtractionResult

NEAR_WHITE_THRESHOLD = 230
TINY_FONT_PT = 3.0
OFF_SLIDE_MARGIN_EMU = Emu(0.1 * 914400)  # ~0.1 pulgada de tolerancia


def _rgb_is_near_white(rgb) -> bool:
    try:
        r, g, b = rgb[0], rgb[1], rgb[2]
    except Exception:
        return False
    return r >= NEAR_WHITE_THRESHOLD and g >= NEAR_WHITE_THRESHOLD and b >= NEAR_WHITE_THRESHOLD


class PptxExtractor(BaseExtractor):
    extensions = (".pptx",)

    def extract(self, path: str) -> ExtractionResult:
        result = ExtractionResult()
        prs = Presentation(path)
        sw, sh = prs.slide_width, prs.slide_height

        self._metadata(prs, result)

        for idx, slide in enumerate(prs.slides, start=1):
            slide_hidden = slide._element.get("show") == "0"
            if slide_hidden:
                result.structural_findings.append(
                    Finding(
                        technique=Technique.HIDDEN_TEXT,
                        severity=Severity.MEDIUM,
                        title=f"Slide oculto (#{idx})",
                        location=f"slide {idx}",
                        evidence='show="0"',
                        detail="El slide esta marcado como oculto; no se muestra al presentar pero su texto se extrae.",
                        hidden=True,
                    )
                )

            for shape in slide.shapes:
                # texto alternativo
                alt = getattr(shape, "alternative_text", "") or ""
                if alt.strip():
                    result.spans.append(
                        TextSpan(text=alt.strip(), visible=False,
                                 location=f"slide {idx} (alt-text)", reason=HideReason.ALT_TEXT)
                    )

                if not shape.has_text_frame:
                    continue

                off_slide = self._is_off_slide(shape, sw, sh)

                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        text = run.text
                        if not text or not text.strip():
                            continue

                        reason = None
                        size_pt = None
                        if slide_hidden:
                            reason = HideReason.NOTES  # contexto: venia de un slide oculto
                        elif off_slide:
                            reason = HideReason.OFF_PAGE
                        else:
                            # fuente diminuta
                            if run.font.size is not None:
                                size_pt = run.font.size.pt
                                if size_pt < TINY_FONT_PT:
                                    reason = HideReason.TINY_FONT
                            # color casi blanco (puede lanzar si es color de tema)
                            if reason is None:
                                try:
                                    if run.font.color and run.font.color.type is not None:
                                        if _rgb_is_near_white(run.font.color.rgb):
                                            reason = HideReason.NEAR_WHITE
                                except Exception:
                                    pass

                        result.spans.append(
                            TextSpan(text=text, visible=reason is None,
                                     location=f"slide {idx}", reason=reason, font_size=size_pt)
                        )

            # notas del orador
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text if slide.notes_slide.notes_text_frame else ""
                if notes and notes.strip():
                    result.spans.append(
                        TextSpan(text=notes.strip(), visible=False,
                                 location=f"slide {idx} (speaker notes)", reason=HideReason.NOTES)
                    )

        return result

    @staticmethod
    def _is_off_slide(shape, sw, sh) -> bool:
        try:
            left, top, w, h = shape.left, shape.top, shape.width, shape.height
        except Exception:
            return False
        if left is None or top is None:
            return False
        w = w or 0
        h = h or 0
        m = OFF_SLIDE_MARGIN_EMU
        return (
            left + w < -m
            or top + h < -m
            or left > sw + m
            or top > sh + m
        )

    def _metadata(self, prs, result: ExtractionResult) -> None:
        props = prs.core_properties
        meta = {}
        for attr in ("title", "subject", "author", "keywords", "comments", "category", "last_modified_by"):
            val = getattr(props, attr, None)
            if val and str(val).strip():
                meta[attr] = str(val)
                result.spans.append(
                    TextSpan(text=str(val), visible=False,
                             location=f"metadata:{attr}", reason=HideReason.METADATA)
                )
        result.metadata = meta
